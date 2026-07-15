"""Tests for the unified ReadFileTool (doyoutrade.tools.file_tools.ReadFileTool).

After the refactor, ``read_file`` is the single in-process file reader:
  - Text-like files (.py, .txt, .json, .yaml, .csv, etc.) → line-numbered content.
  - Document files (.pdf, .docx, .pptx, .xlsx) → extracted plain text, no line numbers.
  - Images (.png, .jpg, etc.) → OCR text (requires rapidocr) or error_code=image_ocr_unavailable.
  - Unknown binary → error_code=unsupported_file_type.
  - No sandbox enforcement on read.
"""
from __future__ import annotations

import json
import os
import tempfile
import unittest
from pathlib import Path

from doyoutrade.tools.file_tools import ReadFileTool
from doyoutrade.tools import build_default_tool_registry


class TestReadTxtFile(unittest.TestCase):
    def setUp(self):
        self.tool = ReadFileTool()

    def test_read_txt_file(self):
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as f:
            f.write("Hello, world!\nLine two\nLine three")
            path = f.name
        try:
            result = self.tool.execute(file_path=path)
            self.assertEqual(result["status"], "ok")
            self.assertEqual(result["format"], "plain")
            self.assertIn("Hello, world!", result["content"])
        finally:
            os.unlink(path)

    def test_read_txt_has_line_numbers(self):
        """Text files must return content with 1-indexed line-number prefix."""
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as f:
            f.write("first\nsecond\nthird\n")
            path = f.name
        try:
            result = self.tool.execute(file_path=path)
            self.assertEqual(result["status"], "ok")
            self.assertIn("1\tfirst", result["content"])
            self.assertIn("2\tsecond", result["content"])
            self.assertIn("3\tthird", result["content"])
        finally:
            os.unlink(path)

    def test_read_python_file_has_line_numbers(self):
        """Python files are line-numbered text."""
        with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False) as f:
            f.write('"""Module docstring."""\n\ndef hello():\n    return "world"\n')
            path = f.name
        try:
            result = self.tool.execute(file_path=path)
            self.assertEqual(result["status"], "ok")
            self.assertEqual(result["format"], "plain")
            # Line 1 should be the docstring line
            self.assertIn("1\t", result["content"])
            self.assertIn("Module docstring", result["content"])
            self.assertIn("def hello", result["content"])
        finally:
            os.unlink(path)

    def test_read_csv(self):
        with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False) as f:
            f.write("name,age,city\nAlice,30,Beijing\nBob,25,Shanghai")
            path = f.name
        try:
            result = self.tool.execute(file_path=path)
            self.assertEqual(result["status"], "ok")
            self.assertEqual(result["format"], "plain")
            self.assertIn("Alice", result["content"])
            self.assertIn("Beijing", result["content"])
        finally:
            os.unlink(path)

    def test_read_json(self):
        with tempfile.NamedTemporaryFile(mode="w", suffix=".json", delete=False) as f:
            json.dump({"key": "value", "number": 42, "list": [1, 2, 3]}, f)
            path = f.name
        try:
            result = self.tool.execute(file_path=path)
            self.assertEqual(result["status"], "ok")
            self.assertEqual(result["format"], "plain")
            self.assertIn("key", result["content"])
            self.assertIn("value", result["content"])
        finally:
            os.unlink(path)

    def test_read_yaml_file(self):
        with tempfile.NamedTemporaryFile(mode="w", suffix=".yaml", delete=False) as f:
            f.write("name: test\nversion: 1.0\nitems:\n  - a\n  - b\n")
            path = f.name
        try:
            result = self.tool.execute(file_path=path)
            self.assertEqual(result["status"], "ok")
            self.assertEqual(result["format"], "plain")
            self.assertIn("name", result["content"])
            self.assertIn("test", result["content"])
        finally:
            os.unlink(path)

    def test_read_nonexistent_file(self):
        result = self.tool.execute(file_path="/nonexistent/path/to/file.txt")
        self.assertEqual(result["status"], "error")
        self.assertEqual(result["error_code"], "file_not_found")

    def test_read_relative_path_rejected(self):
        """Relative paths must be rejected with invalid_path."""
        result = self.tool.execute(file_path="relative/path.txt")
        self.assertEqual(result["status"], "error")
        self.assertEqual(result["error_code"], "invalid_path")

    def test_truncation(self):
        large_content = "A" * 150000
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as f:
            f.write(large_content)
            path = f.name
        try:
            result = self.tool.execute(file_path=path)
            self.assertEqual(result["status"], "ok")
            self.assertTrue(result["truncated"])
            # After line-numbering the truncated content will be > 100_000 bytes
            # due to prefix; check it's truncated and content is long
            self.assertGreater(len(result["content"]), 0)
        finally:
            os.unlink(path)

    def test_encoding_fallback(self):
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
            path = f.name
        try:
            Path(path).write_bytes("你好世界".encode("gbk"))
            result = self.tool.execute(file_path=path)
            self.assertEqual(result["status"], "ok")
            self.assertIn("你好世界", result["content"])
        finally:
            os.unlink(path)

    def test_read_outside_any_sandbox_succeeds(self):
        """read_file works on any path — no sandbox restriction."""
        with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False) as f:
            f.write("x = 42\n")
            path = f.name
        try:
            # No sandbox registered for /tmp — should still succeed.
            result = self.tool.execute(file_path=path)
            self.assertEqual(result["status"], "ok")
            self.assertIn("x = 42", result["content"])
        finally:
            os.unlink(path)

    def test_offset_and_limit_text(self):
        """offset/limit (byte-based) apply to text files."""
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as f:
            f.write("a" * 1000)
            path = f.name
        try:
            result = self.tool.execute(file_path=path, offset=100, limit=200)
            self.assertEqual(result["status"], "ok")
            self.assertEqual(result["char_count"], 200)
        finally:
            os.unlink(path)


class TestReadDocxFile(unittest.TestCase):
    def setUp(self):
        self.tool = ReadFileTool()

    def test_read_docx(self):
        try:
            from docx import Document
        except ImportError:
            self.skipTest("python-docx not installed")

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc_path = f.name

        try:
            doc = Document()
            doc.add_paragraph("Hello from Word")
            doc.add_paragraph("Second paragraph")
            doc.save(doc_path)

            result = self.tool.execute(file_path=doc_path)
            self.assertEqual(result["status"], "ok")
            self.assertEqual(result["format"], "docx")
            self.assertIn("Hello from Word", result["content"])
            self.assertIn("Second paragraph", result["content"])
            # Document extraction should NOT have line-number prefix
            self.assertNotIn("1\tHello", result["content"])
        finally:
            os.unlink(doc_path)


class TestReadXlsxFile(unittest.TestCase):
    def setUp(self):
        self.tool = ReadFileTool()

    def test_read_xlsx(self):
        try:
            import openpyxl
        except ImportError:
            self.skipTest("openpyxl not installed")

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            xlsx_path = f.name

        try:
            wb = __import__("openpyxl").Workbook()
            ws = wb.active
            ws.title = "Sheet1"
            ws["A1"] = "Name"
            ws["B1"] = "Value"
            ws["A2"] = "Foo"
            ws["B2"] = 123
            wb.save(xlsx_path)

            result = self.tool.execute(file_path=xlsx_path)
            self.assertEqual(result["status"], "ok")
            self.assertEqual(result["format"], "xlsx")
            self.assertIn("Name", result["content"])
            self.assertIn("Foo", result["content"])
            self.assertIn("123", result["content"])
            # Document extraction should NOT have line-number prefix
            self.assertNotIn("1\tName", result["content"])
        finally:
            os.unlink(xlsx_path)

    def test_read_xls_returns_error(self):
        """Old .xls format returns a descriptive error."""
        with tempfile.NamedTemporaryFile(suffix=".xls", delete=False) as f:
            xls_path = f.name

        try:
            # Write minimal content so the file exists
            Path(xls_path).write_text("fake xls content")
            result = self.tool.execute(file_path=xls_path)
            self.assertEqual(result["status"], "error")
            self.assertEqual(result["error_code"], "old_xls_format")
        finally:
            os.unlink(xls_path)


class TestReadPptxFile(unittest.TestCase):
    def setUp(self):
        self.tool = ReadFileTool()

    def test_read_pptx(self):
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx not installed")

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name

        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Slide Title"
            slide.placeholders[1].text = "Slide content text"
            prs.save(pptx_path)

            result = self.tool.execute(file_path=pptx_path)
            self.assertEqual(result["status"], "ok")
            self.assertEqual(result["format"], "pptx")
            self.assertIn("Slide Title", result["content"])
            self.assertIn("Slide content text", result["content"])
            # No line numbers for document types
            self.assertNotIn("1\tSlide", result["content"])
        finally:
            os.unlink(pptx_path)


class TestReadPdfFile(unittest.TestCase):
    def setUp(self):
        self.tool = ReadFileTool()

    def test_read_pdf_no_line_numbers(self):
        """PDF extraction returns plain text WITHOUT line-number prefix."""
        try:
            import pypdfium2
        except ImportError:
            self.skipTest("pypdfium2 not installed")

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            pdf_path = f.name

        try:
            pdf_bytes = self._create_minimal_pdf(b"Hello PDF content on page one.")
            with open(pdf_path, "wb") as f:
                f.write(pdf_bytes)

            result = self.tool.execute(file_path=pdf_path)
            self.assertEqual(result["status"], "ok")
            self.assertEqual(result["format"], "pdf")
            self.assertIn("Hello PDF content", result["content"])
            self.assertGreater(result["total_pages"], 0)
            self.assertEqual(result["pages_read"], result["total_pages"])
            # No line numbers for PDF
            self.assertNotIn("1\t", result["content"])
        finally:
            os.unlink(pdf_path)

    def _create_minimal_pdf(self, content_bytes: bytes) -> bytes:
        text_stream = (
            b"BT\n"
            b"/F1 12 Tf\n"
            b"100 700 Td\n"
            b"(" + content_bytes + b") Tj\n"
            b"ET"
        )
        body = (
            b"<< /Length " + str(len(text_stream)).encode() + b" >>\n"
            b"stream\n" + text_stream + b"\nendstream"
        )
        return (
            b"%PDF-1.4\n"
            b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
            b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
            b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"
            b"4 0 obj\n" + body + b"\nendobj\n"
            b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
            b"xref\n0 6\n"
            b"0000000000 65535 f \n"
            b"0000000009 00000 n \n"
            b"0000000058 00000 n \n"
            b"0000000115 00000 n \n"
            b"0000000206 00000 n \n"
            b"0000000385 00000 n \n"
            b"trailer\n<< /Size 6 /Root 1 0 R >>\n"
            b"startxref\n" + str(385 + len(body)).encode() + b"\n"
            b"%%EOF"
        )


class TestReadImageFile(unittest.TestCase):
    def setUp(self):
        self.tool = ReadFileTool()

    def test_read_image_returns_expected_result(self):
        """Image files are dispatched to OCR handler.

        When rapidocr is not installed → error_code=image_ocr_unavailable.
        When rapidocr is installed but the image is unreadable → error_code=io_error.
        When rapidocr is installed and the image is readable → status=ok, format=image.
        Any of these outcomes is acceptable; what matters is that the format
        dispatch reached the image handler (not the text handler or unsupported_file_type).
        """
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            png_path = f.name

        try:
            png_bytes = self._create_minimal_png()
            with open(png_path, "wb") as fh:
                fh.write(png_bytes)

            result = self.tool.execute(file_path=png_path)
            # The dispatch must NOT return unsupported_file_type for .png
            self.assertNotEqual(result.get("error_code"), "unsupported_file_type")
            # Acceptable outcomes:
            #   - rapidocr not installed: image_ocr_unavailable
            #   - rapidocr installed, image bad: io_error
            #   - rapidocr installed, image ok: status=ok, format=image
            if result["status"] == "ok":
                self.assertEqual(result["format"], "image")
            else:
                self.assertIn(
                    result.get("error_code"),
                    {"image_ocr_unavailable", "io_error"},
                    f"Unexpected error_code for image: {result}",
                )
        finally:
            os.unlink(png_path)

    def _create_minimal_png(self) -> bytes:
        import struct, zlib

        def png_chunk(chunk_type: bytes, data: bytes) -> bytes:
            chunk = chunk_type + data
            return struct.pack(">I", len(data)) + chunk + struct.pack(">I", zlib.crc32(chunk) & 0xFFFFFFFF)

        ihdr = struct.pack(">IIBBBBB", 8, 8, 8, 0, 0, 0, 0)
        idat_data = zlib.compress(b"\x00" + b"\x00" * 64, 9)
        return b"\x89PNG\r\n\x1a\n" + png_chunk(b"IHDR", ihdr) + png_chunk(b"IDAT", idat_data) + png_chunk(b"IEND", b"")


class TestToolRegistration(unittest.TestCase):
    def test_registered_in_default_registry(self):
        registry = build_default_tool_registry()
        tool_defs = registry.definitions()
        names = [d["function"]["name"] for d in tool_defs]
        # read_upload_file has been removed; read_file is the unified reader.
        self.assertNotIn("read_upload_file", names)
        self.assertIn("read_file", names)

    def test_read_file_schema(self):
        tool = ReadFileTool()
        schema = tool.to_openai_schema()
        self.assertEqual(schema["type"], "function")
        self.assertEqual(schema["function"]["name"], "read_file")
        self.assertIn("file_path", schema["function"]["parameters"]["properties"])
        self.assertIn("file_path", schema["function"]["parameters"]["required"])
        # offset and limit are optional extras
        self.assertIn("offset", schema["function"]["parameters"]["properties"])
        self.assertIn("limit", schema["function"]["parameters"]["properties"])


if __name__ == "__main__":
    unittest.main()
