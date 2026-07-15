"""Claude Code-style file primitives (in-process agent tools).

These tools mirror Read / Write / Edit / Glob from the host CLI so models
can reuse their existing fluency. They operate on **absolute ``file_path``
arguments**.

Sandbox enforcement:
  - ``read_file`` / ``list_files``: **no sandbox** — reading never mutates;
    ``execute_bash`` already lets the agent read anywhere, so enforcing a
    sandbox here was only friction without security benefit.
  - ``write_file`` / ``edit_file``: **sandbox-enforced** via the
    ``_sandbox`` module-level registry. Two kinds of roots are accepted:
    (1) strategy-authoring ``work_dir`` roots, registered on
    ``open_strategy_authoring`` and unregistered on cancel / finalize; and
    (2) the standing private knowledge base ``~/.doyoutrade/knowledge``,
    registered permanently at tool-registry build time
    (``register_knowledge_sandbox``). Writes to the knowledge base are gated
    behaviourally (only when the user explicitly asks) by the main-agent
    prompt + the ``doyoutrade-knowledge`` skill, not by the sandbox.

No ``session_id`` parameter is needed: the registry tracks which directories
are currently open, and write/edit tools reject any path that does not live
inside one of them.

These tools are NOT subclasses of ``OperationHandler`` — they return plain
``dict`` synchronously — but they ARE registered on the agent's tool surface
(in ``doyoutrade/tools/__init__.py::build_default_tool_registry``) alongside
the framework primitives. The async dispatchers handle the sync-dict return
via ``adapt_sync_dict_to_tool_result``.

Debug-event emission is handled by the dispatcher boundary, not by these
tools individually.

Error codes (stable tokens referenced in skill docs):
  - ``invalid_path``              — file_path is not an absolute path (read)
  - ``file_not_found``            — target does not exist (read / edit)
  - ``path_outside_workspace``    — file_path not inside any active sandbox (write / edit)
  - ``io_error``                  — OS-level IO failure
  - ``old_string_not_found``      — no match for old_string (edit)
  - ``old_string_not_unique``     — old_string matches > 1 time (edit)
  - ``no_op_edit``                — old_string == new_string
  - ``pdf_dep_missing``           — pypdfium2 not installed (read .pdf)
  - ``docx_dep_missing``          — python-docx not installed (read .docx)
  - ``pptx_dep_missing``          — python-pptx not installed (read .pptx)
  - ``xlsx_dep_missing``          — openpyxl not installed (read .xlsx/.xls)
  - ``image_ocr_unavailable``     — rapidocr not installed (read image)
  - ``old_xls_format``            — .xls files are not supported (use .xlsx)
"""
from __future__ import annotations

import json as _json
from dataclasses import dataclass
from pathlib import Path
from typing import Any, ClassVar

from doyoutrade.persistence.strategy_storage import SandboxViolation
from doyoutrade.tools._contract import ContractResult, enforce_kwargs_contract
from doyoutrade.tools._sandbox import resolve_path as _resolve_sandbox_path

MAX_TEXT_LENGTH = 100_000

# Extensions treated as plain text (line-numbered output).
_TEXT_EXTENSIONS = frozenset({
    ".py", ".txt", ".md", ".json", ".yaml", ".yml", ".csv",
    ".html", ".xml", ".sh", ".toml", ".ini", ".ts", ".tsx",
    ".js", ".jsx", ".css", ".scss", ".sql", ".rb", ".go",
    ".java", ".c", ".cpp", ".h", ".rs", ".swift", ".kt",
    ".r", ".m", ".lua", ".pl", ".php",
})
_DOCUMENT_EXTENSIONS = frozenset({".pdf", ".docx", ".pptx", ".xlsx", ".xls"})
_IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".tiff"})

# Shared sentinel for empty maps; each tool declares its own.
_EMPTY_LIFTS: dict[str, Any] = {}
_EMPTY_SUGGESTED: dict[str, str] = {}


@dataclass
class _FileToolBase:
    """Shared base for the sandboxed / unrestricted file primitives.

    No external dependencies are injected — write/edit tools read from
    the module-level ``_sandbox`` registry instead of receiving a callback.
    """

    # Subclasses declare the set of allowed top-level kwargs via parameters.
    _allowed: ClassVar[frozenset[str]] = frozenset()

    # Empty sentinel; concrete subclasses override with their own schema.
    parameters: ClassVar[dict] = {
        "type": "object",
        "additionalProperties": False,
        "properties": {},
        "required": [],
    }

    def to_openai_schema(self) -> dict[str, Any]:
        """Return an OpenAI-compatible function schema so these tools can be
        registered in ``OperationRegistry`` alongside ``OperationHandler``
        subclasses.
        """
        return {
            "type": "function",
            "function": {
                "name": self.name,
                "description": self.description,
                "parameters": self.parameters,
            },
        }

    def _contract(self, kwargs: dict[str, Any]) -> ContractResult:
        props = self.parameters.get("properties") or {}
        allowed = frozenset(props.keys()) if props else self._allowed
        return enforce_kwargs_contract(
            kwargs,
            allowed_top_level=allowed,
            suggested_paths=_EMPTY_SUGGESTED,
            legacy_lifts=_EMPTY_LIFTS,
        )

    def _resolve_sandbox(self, file_path: str) -> Path:
        """Resolve with sandbox enforcement (for write/edit tools)."""
        return _resolve_sandbox_path(file_path)

    def _resolve_unrestricted(self, file_path: str) -> Path:
        """Resolve without sandbox enforcement (for read tools)."""
        return Path(file_path).resolve(strict=False)


# ---------------------------------------------------------------------------
# Document / image extraction helpers
# ---------------------------------------------------------------------------

def _extract_pdf(path: Path) -> dict[str, Any]:
    """Extract text from a PDF file using pypdfium2."""
    try:
        import pypdfium2
    except ImportError:
        return {"status": "error", "error_code": "pdf_dep_missing",
                "message": "pypdfium2 not installed"}
    try:
        with pypdfium2.PdfDocument(str(path)) as pdf:
            total_pages = len(pdf)
            lines = []
            for page_num in range(total_pages):
                page = pdf[page_num]
                textpage = page.get_textpage()
                page_text = textpage.get_text_bounded()
                if page_text.strip():
                    lines.append(f"--- Page {page_num + 1} ---\n{page_text}")
            text = "\n".join(lines)
    except Exception as exc:
        return {"status": "error", "error_code": "io_error",
                "message": f"PDF read failed: {exc}"}
    truncated, text = _truncate(text)
    return {
        "status": "ok",
        "format": "pdf",
        "text": text,
        "truncated": truncated,
        "total_pages": total_pages,
        "pages_read": total_pages,
    }


def _extract_docx(path: Path) -> dict[str, Any]:
    """Extract text from a .docx file using python-docx."""
    try:
        import docx
    except ImportError:
        return {"status": "error", "error_code": "docx_dep_missing",
                "message": "python-docx not installed"}
    try:
        doc = docx.Document(str(path))
        lines = [para.text for para in doc.paragraphs if para.text.strip()]
        text = "\n".join(lines)
    except Exception as exc:
        return {"status": "error", "error_code": "io_error",
                "message": f"DOCX read failed: {exc}"}
    truncated, text = _truncate(text)
    return {"status": "ok", "format": "docx", "text": text, "truncated": truncated}


def _extract_xlsx(path: Path, ext: str) -> dict[str, Any]:
    """Extract text from a .xlsx file using openpyxl."""
    if ext == ".xls":
        return {"status": "error", "error_code": "old_xls_format",
                "message": "old .xls format is not supported; save as .xlsx first"}
    try:
        import openpyxl
    except ImportError:
        return {"status": "error", "error_code": "xlsx_dep_missing",
                "message": "openpyxl not installed"}
    try:
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        ws = wb.active
        lines = []
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                lines.append(",".join(str(c) if c is not None else "" for c in row))
        text = "\n".join(lines)
        wb.close()
    except Exception as exc:
        return {"status": "error", "error_code": "io_error",
                "message": f"XLSX read failed: {exc}"}
    truncated, text = _truncate(text)
    return {"status": "ok", "format": "xlsx", "text": text, "truncated": truncated}


def _extract_pptx(path: Path) -> dict[str, Any]:
    """Extract text from a .pptx file using python-pptx."""
    try:
        import pptx
    except ImportError:
        return {"status": "error", "error_code": "pptx_dep_missing",
                "message": "python-pptx not installed"}
    try:
        prs = pptx.Presentation(str(path))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text)
        text = "\n".join(lines)
    except Exception as exc:
        return {"status": "error", "error_code": "io_error",
                "message": f"PPTX read failed: {exc}"}
    truncated, text = _truncate(text)
    return {"status": "ok", "format": "pptx", "text": text, "truncated": truncated}


def _extract_image(path: Path) -> dict[str, Any]:
    """Extract text from an image via OCR (rapidocr)."""
    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        return {"status": "error", "error_code": "image_ocr_unavailable",
                "message": "rapidocr_onnxruntime not installed; OCR unavailable"}
    try:
        ocr = RapidOCR()
        result, elapse = ocr(str(path))
        if result is None:
            text = ""
        else:
            text = "\n".join(line[1] for line in result)
    except Exception as exc:
        return {"status": "error", "error_code": "io_error",
                "message": f"image OCR failed: {exc}"}
    truncated, text = _truncate(text)
    return {"status": "ok", "format": "image", "text": text, "truncated": truncated}


def _truncate(text: str) -> tuple[bool, str]:
    if len(text) > MAX_TEXT_LENGTH:
        return True, text[:MAX_TEXT_LENGTH]
    return False, text


def _read_text(path: Path, offset: int = 0, limit: int = MAX_TEXT_LENGTH) -> dict[str, Any]:
    """Read a text file with encoding fallback and optional byte-based offset/limit.

    Returns content with 1-indexed line-number prefix (``"{n}\\t{line}"``)
    applied AFTER slicing. offset/limit are byte-based for backward compat.
    """
    encodings = ["utf-8", "gbk", "gb2312", "big5", "latin-1"]
    raw = ""
    for enc in encodings:
        try:
            raw = path.read_text(encoding=enc)
            break
        except UnicodeDecodeError:
            if enc == encodings[-1]:
                raw = path.read_text(encoding="latin-1", errors="replace")
            continue
        except OSError as exc:
            return {"status": "error", "error_code": "io_error",
                    "message": f"read failed: {exc}"}

    total_len = len(raw)
    if offset >= total_len:
        sliced = ""
    else:
        sliced = raw[offset: offset + limit]

    truncated_by_limit = (len(sliced) == limit) and (offset + limit < total_len)
    over_max, sliced = _truncate(sliced)
    truncated = truncated_by_limit or over_max

    numbered = "\n".join(
        f"{i + 1}\t{line}" for i, line in enumerate(sliced.splitlines())
    )
    return {
        "status": "ok",
        "format": "plain",
        "content": numbered,
        "char_count": len(sliced),
        "truncated": truncated,
    }


def _dispatch_read(path: Path, offset: int = 0, limit: int = MAX_TEXT_LENGTH) -> dict[str, Any]:
    """Dispatch to the correct extractor based on file extension.

    Text-like files → line-numbered content under key ``content``.
    Document/image files → plain text under key ``content`` (no line numbers).
    """
    ext = path.suffix.lower()

    if ext in _TEXT_EXTENSIONS:
        return _read_text(path, offset=offset, limit=limit)
    elif ext == ".pdf":
        result = _extract_pdf(path)
        if result["status"] == "ok":
            result["content"] = result.pop("text")
        return result
    elif ext == ".docx":
        result = _extract_docx(path)
        if result["status"] == "ok":
            result["content"] = result.pop("text")
        return result
    elif ext in {".xlsx", ".xls"}:
        result = _extract_xlsx(path, ext)
        if result["status"] == "ok":
            result["content"] = result.pop("text")
        return result
    elif ext == ".pptx":
        result = _extract_pptx(path)
        if result["status"] == "ok":
            result["content"] = result.pop("text")
        return result
    elif ext in _IMAGE_EXTENSIONS:
        result = _extract_image(path)
        if result["status"] == "ok":
            result["content"] = result.pop("text")
        return result
    else:
        # Unknown extension — try reading as UTF-8 text; fall back gracefully.
        try:
            raw = path.read_bytes()
            text = raw.decode("utf-8", errors="strict")
            result = _read_text(path, offset=offset, limit=limit)
            return result
        except (UnicodeDecodeError, OSError):
            pass
        # Binary file we can't decode — report as unsupported.
        return {
            "status": "error",
            "error_code": "unsupported_file_type",
            "message": f"Unsupported file type: {ext}",
        }


# ---------------------------------------------------------------------------
# ReadFileTool — unrestricted, multimodal dispatch
# ---------------------------------------------------------------------------

@dataclass
class ReadFileTool(_FileToolBase):
    """Read any file: source code is returned with 1-indexed line numbers;
    documents (PDF / DOCX / PPTX / XLSX) and images return extracted plain
    text without line numbers.

    **No sandbox enforcement** — reading is non-mutating and the agent can
    already read anywhere via ``execute_bash``.

    ``error_code`` tokens: ``invalid_path``, ``file_not_found``, ``io_error``,
    ``pdf_dep_missing``, ``docx_dep_missing``, ``pptx_dep_missing``,
    ``xlsx_dep_missing``, ``image_ocr_unavailable``, ``unsupported_file_type``.

    Optional ``offset`` / ``limit`` are byte-based and apply to text files
    only (ignored for binary documents).
    """

    name: str = "read_file"
    description: str = (
        "Read any file at an absolute path. Source code / text files are "
        "returned with 1-indexed line numbers. Documents (PDF, DOCX, PPTX, "
        "XLSX) and images return extracted plain text. No sandbox restriction "
        "— you can read any path the OS allows."
    )
    parameters: ClassVar[dict] = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Absolute path to the file to read.",
            },
            "offset": {
                "type": "integer",
                "description": (
                    "Byte offset to start reading from (text files only). "
                    "Defaults to 0."
                ),
                "default": 0,
            },
            "limit": {
                "type": "integer",
                "description": (
                    "Maximum bytes to read after offset (text files only). "
                    f"Defaults to {MAX_TEXT_LENGTH}."
                ),
                "default": MAX_TEXT_LENGTH,
            },
        },
        "required": ["file_path"],
    }
    _allowed: ClassVar[frozenset[str]] = frozenset({"file_path", "offset", "limit"})

    def execute(self, **kwargs: Any) -> dict[str, Any]:
        contract = self._contract(kwargs)
        if contract.error is not None:
            return {
                "status": "error",
                "error_code": contract.error.get("type", "validation_error"),
                "message": contract.error.get("message", "invalid arguments"),
            }
        file_path: str = contract.kwargs["file_path"]
        offset: int = int(contract.kwargs.get("offset") or 0)
        limit: int = int(contract.kwargs.get("limit") or MAX_TEXT_LENGTH)

        if not Path(file_path).is_absolute():
            return {
                "status": "error",
                "error_code": "invalid_path",
                "message": f"file_path must be an absolute path, got: {file_path!r}",
            }

        target = self._resolve_unrestricted(file_path)

        if not target.is_file():
            return {
                "status": "error",
                "error_code": "file_not_found",
                "message": f"{file_path} does not exist",
            }

        result = _dispatch_read(target, offset=offset, limit=limit)
        if result.get("status") == "ok":
            result["file_path"] = file_path
        return result


@dataclass
class WriteFileTool(_FileToolBase):
    """Write (overwrite) a file inside an active sandbox root.

    Creates parent directories as needed.  ``error_code`` tokens:
    ``path_outside_workspace``, ``io_error``.
    """

    name: str = "write_file"
    description: str = (
        "Write (overwrite) a file inside an active sandbox root: either the open "
        "strategy-authoring work_dir, or the private knowledge base "
        "~/.doyoutrade/knowledge (write there only when the user explicitly asks). "
        "file_path must be an absolute path within one of those roots."
    )
    parameters: ClassVar[dict] = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Absolute path to the file inside the active sandbox.",
            },
            "content": {
                "type": "string",
                "description": "New file contents. Parent directories are created as needed.",
            },
        },
        "required": ["file_path", "content"],
    }
    _allowed: ClassVar[frozenset[str]] = frozenset({"file_path", "content"})

    def execute(self, **kwargs: Any) -> dict[str, Any]:
        contract = self._contract(kwargs)
        if contract.error is not None:
            return {
                "status": "error",
                "error_code": contract.error.get("type", "validation_error"),
                "message": contract.error.get("message", "invalid arguments"),
            }
        file_path: str = contract.kwargs["file_path"]
        content: str = contract.kwargs["content"]

        try:
            target = self._resolve_sandbox(file_path)
        except SandboxViolation as exc:
            return {
                "status": "error",
                "error_code": "path_outside_workspace",
                "message": str(exc),
            }

        try:
            target.parent.mkdir(parents=True, exist_ok=True)
            target.write_text(content)
        except OSError as exc:
            return {
                "status": "error",
                "error_code": "io_error",
                "message": f"write failed: {exc}",
            }

        return {
            "status": "ok",
            "file_path": file_path,
            "bytes_written": len(content.encode("utf-8")),
        }


@dataclass
class EditFileTool(_FileToolBase):
    """Replace ``old_string`` with ``new_string`` in a sandbox file.

    ``old_string`` must appear exactly once unless ``replace_all=True``.
    ``error_code`` tokens: ``path_outside_workspace``, ``file_not_found``,
    ``no_op_edit``, ``old_string_not_found``, ``old_string_not_unique``,
    ``io_error``.
    """

    name: str = "edit_file"
    description: str = (
        "Replace ``old_string`` with ``new_string`` in a sandbox file (strategy "
        "authoring work_dir, or ~/.doyoutrade/knowledge when the user asks). "
        "``old_string`` must be unique unless ``replace_all=true``."
    )
    parameters: ClassVar[dict] = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Absolute path to the file inside the active sandbox.",
            },
            "old_string": {
                "type": "string",
                "description": "Substring to replace. Must occur exactly once unless replace_all=true.",
            },
            "new_string": {
                "type": "string",
                "description": "Replacement substring.",
            },
            "replace_all": {
                "type": "boolean",
                "default": False,
                "description": "If true, replace every occurrence.",
            },
        },
        "required": ["file_path", "old_string", "new_string"],
    }
    _allowed: ClassVar[frozenset[str]] = frozenset(
        {"file_path", "old_string", "new_string", "replace_all"}
    )

    def execute(self, **kwargs: Any) -> dict[str, Any]:
        contract = self._contract(kwargs)
        if contract.error is not None:
            return {
                "status": "error",
                "error_code": contract.error.get("type", "validation_error"),
                "message": contract.error.get("message", "invalid arguments"),
            }
        file_path: str = contract.kwargs["file_path"]
        old_string: str = contract.kwargs["old_string"]
        new_string: str = contract.kwargs["new_string"]
        replace_all: bool = contract.kwargs.get("replace_all", False)

        if old_string == new_string:
            return {
                "status": "error",
                "error_code": "no_op_edit",
                "message": "old_string and new_string are identical",
            }

        try:
            target = self._resolve_sandbox(file_path)
        except SandboxViolation as exc:
            return {
                "status": "error",
                "error_code": "path_outside_workspace",
                "message": str(exc),
            }

        if not target.is_file():
            return {
                "status": "error",
                "error_code": "file_not_found",
                "message": f"{file_path} does not exist",
            }

        try:
            body = target.read_text()
        except OSError as exc:
            return {
                "status": "error",
                "error_code": "io_error",
                "message": f"read failed: {exc}",
            }

        count = body.count(old_string)

        if count == 0:
            return {
                "status": "error",
                "error_code": "old_string_not_found",
                "message": "no match for old_string",
            }

        if count > 1 and not replace_all:
            return {
                "status": "error",
                "error_code": "old_string_not_unique",
                "message": (
                    f"old_string matches {count} times; pass replace_all=true "
                    "or provide more surrounding context to make it unique"
                ),
            }

        if replace_all:
            new_body = body.replace(old_string, new_string)
            replacements = count
        else:
            new_body = body.replace(old_string, new_string, 1)
            replacements = 1

        try:
            target.write_text(new_body)
        except OSError as exc:
            return {
                "status": "error",
                "error_code": "io_error",
                "message": f"write failed: {exc}",
            }

        return {"status": "ok", "file_path": file_path, "replacements": replacements}


@dataclass
class ListFilesTool(_FileToolBase):
    """List all files under a directory as POSIX relative paths.

    The ``directory`` parameter must be an absolute path. No sandbox
    enforcement — listing is read-only and the agent can already list
    any directory via ``execute_bash``.

    ``error_code`` tokens: ``file_not_found``, ``io_error``.
    """

    name: str = "list_files"
    description: str = (
        "List all files under a directory as POSIX relative paths. "
        "directory must be an absolute path. No sandbox restriction."
    )
    parameters: ClassVar[dict] = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "directory": {
                "type": "string",
                "description": (
                    "Absolute path to the directory to list."
                ),
            },
        },
        "required": ["directory"],
    }
    _allowed: ClassVar[frozenset[str]] = frozenset({"directory"})

    def execute(self, **kwargs: Any) -> dict[str, Any]:
        contract = self._contract(kwargs)
        if contract.error is not None:
            return {
                "status": "error",
                "error_code": contract.error.get("type", "validation_error"),
                "message": contract.error.get("message", "invalid arguments"),
            }
        directory: str = contract.kwargs["directory"]

        root = Path(directory).resolve(strict=False)

        if not root.is_dir():
            return {
                "status": "error",
                "error_code": "file_not_found",
                "message": f"{directory} does not exist or is not a directory",
            }

        try:
            # Walk recursively, collect files only, convert to POSIX relative
            files = sorted(
                p.relative_to(root).as_posix()
                for p in root.rglob("*")
                if p.is_file()
            )
        except OSError as exc:
            return {
                "status": "error",
                "error_code": "io_error",
                "message": f"list failed: {exc}",
            }

        return {"status": "ok", "files": files}


__all__ = [
    "ReadFileTool",
    "WriteFileTool",
    "EditFileTool",
    "ListFilesTool",
]
